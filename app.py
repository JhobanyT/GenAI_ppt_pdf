from flask import Flask, render_template, request, send_file, jsonify
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import os
import openai
from pptx import Presentation
import random
import re
from dotenv import load_dotenv
import subprocess

# Cargar las variables de entorno desde el archivo .env
load_dotenv()

# Obtener la API key desde las variables de entorno
openai.api_key = os.getenv("OPENAI_API_KEY")

# Inicialización de la aplicación Flask
app = Flask(__name__)

# Configuración del limitador de tasa
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["10 per day"]
)

# Prompt para generar la presentación
Prompt = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has a summary.
-At least 8 slides.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END"""

def create_ppt_text(input_text):
    try:
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": Prompt},
                {"role": "user", "content": "The user wants a presentation about " + input_text}
            ],
            temperature=0.5,
        )
        return response.choices[0].message.content
    except Exception as e:
        return str(e)

def create_ppt(text_file, design_number, ppt_name, requested_slide_count):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    first_time = True

    # Eliminar la extensión .txt del nombre del archivo si está presente
    ppt_name = os.path.splitext(ppt_name)[0]

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if slide_count >= requested_slide_count:  # Verificar si alcanzamos el número deseado de diapositivas
                break

            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout para título
                slide.shapes.title.text = header
                slide_count += 1  # Incrementar el contador de diapositivas
                continue

            elif line.startswith('#Slide:'):
                if slide_count > 0:  # Solo si ya hemos creado al menos un slide
                    slide = prs.slides.add_slide(prs.slide_layouts[last_slide_layout_index])
                    slide.shapes.title.text = header
                    body_shape = slide.shapes.placeholders[1]  # Placeholder para el contenido
                    tf = body_shape.text_frame
                    tf.text = content
                    slide_count += 1  # Incrementar el contador de diapositivas
                content = ""  # Reiniciar contenido para la siguiente diapositiva
                last_slide_layout_index = random.choice([1, 7, 8, 9, 10])  # Elegir el layout de forma aleatoria
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    ppt_dir = 'GeneratedPresentations'
    pdf_dir = 'GeneratedPdf'

    # Crear directorios si no existen
    os.makedirs(ppt_dir, exist_ok=True)
    os.makedirs(pdf_dir, exist_ok=True)

    ppt_path = os.path.join(ppt_dir, f'{ppt_name}.pptx')
    prs.save(ppt_path)

    # Verificar que el archivo PPTX existe antes de convertir
    if not os.path.exists(ppt_path):
        raise Exception(f"PPTX file not found: {ppt_path}")

    # Convertir el PPTX a PDF usando LibreOffice
    pdf_path = os.path.join(pdf_dir, f'{ppt_name}.pdf')  # PDF en la carpeta GeneratedPdf
    libreoffice_path = 'C:\\Program Files\\LibreOffice\\program\\soffice.exe'  # Ruta completa
    try:
        result = subprocess.run([libreoffice_path, '--headless', '--convert-to', 'pdf', '--outdir', pdf_dir, ppt_path],
                                capture_output=True, text=True)
        if result.returncode != 0:
            raise Exception(f"LibreOffice conversion error: {result.stderr}")
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
        raise

    # Verificar que el archivo PDF se ha creado
    if not os.path.exists(pdf_path):
        raise Exception(f"PDF file not found: {pdf_path}")

    return {
        "pptx": f"{request.host_url}GeneratedPresentations/{ppt_name}.pptx",
        "pdf": f"{request.host_url}GeneratedPdf/{ppt_name}.pdf"
    }

@app.route('/GeneratedPresentations/<path:path>')
def send_generated_file(path):
    file_path = f'GeneratedPresentations/{path}'
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

@app.route('/GeneratedPdf/<path:path>')
def send_pdf_file(path):
    file_path = f'GeneratedPdf/{path}'
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=False)  # No es como attachment
    return jsonify({"error": "File not found"}), 404

@app.route("/powerpoint")
def powerpoint():
    return render_template("powerpoint.html", charset="utf-8")

@app.route("/")
def home():
    return render_template("powerpoint.html", charset="utf-8")

@app.route("/get")
@limiter.limit("10 per day, key_func=get_remote_address")
def get_bot_response():
    user_text = request.args.get("msg")
    selected_design = request.args.get("design")# Captura el diseño seleccionado
    slide_count = request.args.get("slides")  # Captura el número de diapositivas

    if not user_text:
        return jsonify({"error": "No message provided"}), 400
    
    last_char = user_text[-1]
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', user_text).replace("\n", "")
    design_number = 1

    print(f"Selected Design (raw): {selected_design}")

    if last_char.isdigit():
        design_number = int(last_char)
        input_string = user_text[:-2]
        print("Design Number:", design_number, "selected.")
    else:
        print("No design specified, using default design...")

    if selected_design:  # Verifica que se haya pasado un diseño
        try:
            design_number = int(selected_design)  # Asegúrate de que el diseño sea un número
            print("Selected Design:", design_number)
        except ValueError:
            print("Invalid design selected, using default design...")

    if design_number > 10 or design_number == 0:
        design_number = 1
        print("Unavailable design, using default design...")

    # CANTIDAD PAGINAS

    if slide_count:
        try:
            slide_count = int(slide_count)
            print(f"Slide count received: {slide_count}")
        except ValueError:
            slide_count = 8  # Asigna un valor por defecto si no se puede convertir

    filename_prompt = f"Generate a short, descriptive filename based on the following input: \"{input_string}\". Answer just with the short filename, no other explanation."
    try:
        filename_response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": filename_prompt},
            ],
            temperature=0.5,
            max_tokens=30,
        )
        filename = filename_response.choices[0].message.content.strip().replace(" ", "_").replace(".pptx", "")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    try:
        ppt_text = create_ppt_text(input_string)
        with open(f'Cache/{filename}', 'w', encoding='utf-8') as f:
            f.write(ppt_text)
        paths = create_ppt(f'Cache/{filename}', design_number, filename, slide_count)
        return jsonify(paths)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run()
